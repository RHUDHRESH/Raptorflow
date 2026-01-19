#!/usr/bin/env python3
"""
Comprehensive OCR Testing System
Downloads various content types from internet and tests OCR accuracy
Compares OCR results with known content to verify system performance
"""

import os
import requests
import json
from datetime import datetime
from urllib.parse import urlparse
import time
import asyncio
import aiohttp
from typing import List, Dict, Any, Tuple
import pdf2image
import pytesseract
from PIL import Image
import io

class ComprehensiveOCRTester:
    """Comprehensive OCR testing with knowledge verification"""
    
    def __init__(self):
        self.test_dir = "comprehensive_ocr_test"
        os.makedirs(self.test_dir, exist_ok=True)
        self.setup_ocr()
        
    def setup_ocr(self):
        """Setup OCR configuration"""
        try:
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            print("✅ OCR engine initialized")
        except Exception as e:
            print(f"❌ OCR setup failed: {e}")
    
    async def download_test_content(self) -> List[Dict[str, Any]]:
        """Download diverse content types for OCR testing"""
        
        # Test content with known text/content for verification
        test_content = [
            {
                "url": "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf",
                "category": "pdf_document",
                "description": "W3C Sample PDF Document",
                "expected_content": ["PDF", "document", "sample", "W3C"],
                "filename": "w3c_sample.pdf"
            },
            {
                "url": "https://www.learningcontainer.com/wp-content/uploads/2020/05/sample-mp4-file.mp4",
                "category": "video_file",
                "description": "Sample MP4 Video File",
                "expected_content": ["video", "sample", "mp4"],
                "filename": "sample_video.mp4"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/10/file_example_JPG_100kB.jpg",
                "category": "image_file",
                "description": "Sample JPG Image",
                "expected_content": ["image", "jpg", "sample"],
                "filename": "sample_image.jpg"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/11/file_example_CSV_1000.csv",
                "category": "csv_data",
                "description": "Sample CSV Data File",
                "expected_content": ["csv", "data", "sample"],
                "filename": "sample_data.csv"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/02/file_example_JSON_1kb.json",
                "category": "json_data",
                "description": "Sample JSON Data File",
                "expected_content": ["json", "data", "sample"],
                "filename": "sample_data.json"
            },
            {
                "url": "https://www.learningcontainer.com/wp-content/uploads/2020/04/sample-text-file.txt",
                "category": "text_file",
                "description": "Sample Text File",
                "expected_content": ["text", "file", "sample"],
                "filename": "sample_text.txt"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/08/file_example_XLSX_10.xlsx",
                "category": "excel_file",
                "description": "Sample Excel File",
                "expected_content": ["excel", "xlsx", "sample"],
                "filename": "sample_excel.xlsx"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/08/file_example_PPTX_10kB.pptx",
                "category": "powerpoint_file",
                "description": "Sample PowerPoint File",
                "expected_content": ["powerpoint", "pptx", "sample"],
                "filename": "sample_powerpoint.pptx"
            },
            {
                "url": "https://file-examples.com/storage/fe86ead47062e5b4375c42b5/2017/08/file_example_DOCX_10kB.docx",
                "category": "word_document",
                "description": "Sample Word Document",
                "expected_content": ["word", "docx", "sample"],
                "filename": "sample_word.docx"
            },
            {
                "url": "https://images.unsplash.com/photo-1456513080510-7bf3a84b82f8?w=800&h=600",
                "category": "business_document",
                "description": "Business Document Image",
                "expected_content": ["business", "document", "text"],
                "filename": "business_doc.jpg"
            },
            {
                "url": "https://images.unsplash.com/photo-1586953208448-b95a79798f07?w=800&h=600",
                "category": "spreadsheet_image",
                "description": "Spreadsheet/Financial Data Image",
                "expected_content": ["spreadsheet", "data", "financial"],
                "filename": "spreadsheet.jpg"
            },
            {
                "url": "https://images.unsplash.com/photo-1554224154-2603ffc5d8b2?w=800&h=600",
                "category": "contract_document",
                "description": "Legal Contract Document",
                "expected_content": ["contract", "legal", "document"],
                "filename": "contract.jpg"
            },
            {
                "url": "https://images.unsplash.com/photo-1560185007-c5ca9d2c014d?w=800&h=600",
                "category": "invoice_document",
                "description": "Invoice/Billing Document",
                "expected_content": ["invoice", "billing", "document"],
                "filename": "invoice.jpg"
            },
            {
                "url": "https://images.unsplash.com/photo-1554469384-e58e1667c219?w=800&h=600",
                "category": "report_document",
                "description": "Business Report Document",
                "expected_content": ["report", "business", "analysis"],
                "filename": "report.jpg"
            }
        ]
        
        print("🌐 DOWNLOADING DIVERSE CONTENT FOR OCR TESTING")
        print("=" * 60)
        
        downloaded_files = []
        
        async with aiohttp.ClientSession() as session:
            for i, item in enumerate(test_content):
                try:
                    print(f"\n📥 [{i+1}/{len(test_content)}] Downloading: {item['description']}")
                    print(f"    Category: {item['category']}")
                    print(f"    Expected Content: {', '.join(item['expected_content'])}")
                    
                    headers = {
                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                    }
                    
                    async with session.get(item["url"], headers=headers, timeout=30) as response:
                        if response.status == 200:
                            content = await response.read()
                            
                            # Add timestamp to filename
                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                            filename = f"{item['category']}_{timestamp}_{item['filename']}"
                            filepath = os.path.join(self.test_dir, filename)
                            
                            with open(filepath, 'wb') as f:
                                f.write(content)
                            
                            file_info = {
                                "filename": filename,
                                "filepath": filepath,
                                "size": len(content),
                                "category": item["category"],
                                "description": item["description"],
                                "expected_content": item["expected_content"],
                                "url": item["url"],
                                "content_type": response.headers.get('content-type', 'unknown'),
                                "success": True
                            }
                            
                            downloaded_files.append(file_info)
                            print(f"    ✅ Downloaded: {filename} ({len(content):,} bytes)")
                        else:
                            print(f"    ❌ Failed: HTTP {response.status}")
                            
                except Exception as e:
                    print(f"    ❌ Error: {e}")
                
                # Add delay to avoid rate limiting
                await asyncio.sleep(1)
        
        return downloaded_files
    
    def extract_text_from_file(self, filepath: str, category: str) -> Dict[str, Any]:
        """Extract text using appropriate method based on file type"""
        
        result = {
            "method": "unknown",
            "extracted_text": "",
            "confidence": 0.0,
            "error": None,
            "processing_time": 0.0
        }
        
        start_time = time.time()
        
        try:
            if category == "pdf_document":
                result = self.extract_from_pdf(filepath)
            elif category in ["image_file", "business_document", "spreadsheet_image", "contract_document", "invoice_document", "report_document"]:
                result = self.extract_from_image(filepath)
            elif category == "text_file":
                result = self.extract_from_text_file(filepath)
            elif category == "csv_data":
                result = self.extract_from_csv(filepath)
            elif category == "json_data":
                result = self.extract_from_json(filepath)
            elif category in ["excel_file", "powerpoint_file", "word_document"]:
                result = self.extract_from_office_document(filepath)
            elif category == "video_file":
                result = self.extract_from_video(filepath)
            else:
                result["error"] = f"Unsupported category: {category}"
                
        except Exception as e:
            result["error"] = str(e)
        
        result["processing_time"] = time.time() - start_time
        return result
    
    def extract_from_pdf(self, filepath: str) -> Dict[str, Any]:
        """Extract text from PDF using multiple methods"""
        try:
            # Method 1: Try pdf2image + OCR
            try:
                images = pdf2image.convert_from_path(filepath)
                all_text = ""
                total_confidence = 0
                
                for image in images:
                    text = pytesseract.image_to_string(image, lang='eng')
                    confidence = pytesseract.image_to_data(image, lang='eng', output_type=pytesseract.Output.DICT)
                    avg_conf = sum(confidence['conf']) / len(confidence['conf']) if confidence['conf'] else 0
                    all_text += text + "\n"
                    total_confidence += avg_conf
                
                return {
                    "method": "pdf_ocr",
                    "extracted_text": all_text.strip(),
                    "confidence": total_confidence / len(images) if images else 0,
                    "pages_processed": len(images)
                }
            except Exception as e:
                # Method 2: Try direct PDF text extraction
                try:
                    import PyPDF2
                    with open(filepath, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() + "\n"
                    
                    return {
                        "method": "pdf_direct",
                        "extracted_text": text.strip(),
                        "confidence": 0.8,  # High confidence for direct extraction
                        "pages_processed": len(reader.pages)
                    }
                except:
                    return {"method": "pdf_failed", "error": f"PDF extraction failed: {e}"}
                    
        except Exception as e:
            return {"method": "pdf_error", "error": str(e)}
    
    def extract_from_image(self, filepath: str) -> Dict[str, Any]:
        """Extract text from image using OCR"""
        try:
            image = Image.open(filepath)
            
            # Method 1: Basic OCR
            text = pytesseract.image_to_string(image, lang='eng')
            
            # Get confidence data
            data = pytesseract.image_to_data(image, lang='eng', output_type=pytesseract.Output.DICT)
            confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            return {
                "method": "image_ocr",
                "extracted_text": text.strip(),
                "confidence": avg_confidence / 100,  # Convert to 0-1 scale
                "word_count": len(text.split())
            }
            
        except Exception as e:
            return {"method": "image_ocr_error", "error": str(e)}
    
    def extract_from_text_file(self, filepath: str) -> Dict[str, Any]:
        """Extract text from text file directly"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                text = file.read()
            
            return {
                "method": "text_direct",
                "extracted_text": text.strip(),
                "confidence": 1.0,  # Perfect confidence for direct text
                "character_count": len(text)
            }
            
        except Exception as e:
            return {"method": "text_error", "error": str(e)}
    
    def extract_from_csv(self, filepath: str) -> Dict[str, Any]:
        """Extract text from CSV file"""
        try:
            import pandas as pd
            df = pd.read_csv(filepath)
            text = df.to_string()
            
            return {
                "method": "csv_pandas",
                "extracted_text": text,
                "confidence": 0.9,
                "rows": len(df),
                "columns": len(df.columns)
            }
            
        except Exception as e:
            # Fallback to direct text reading
            try:
                with open(filepath, 'r', encoding='utf-8') as file:
                    text = file.read()
                
                return {
                    "method": "csv_direct",
                    "extracted_text": text,
                    "confidence": 0.8
                }
            except:
                return {"method": "csv_error", "error": str(e)}
    
    def extract_from_json(self, filepath: str) -> Dict[str, Any]:
        """Extract text from JSON file"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert JSON to readable text
            text = json.dumps(data, indent=2)
            
            return {
                "method": "json_direct",
                "extracted_text": text,
                "confidence": 1.0,
                "keys_found": len(str(data).split('"')) // 2
            }
            
        except Exception as e:
            return {"method": "json_error", "error": str(e)}
    
    def extract_from_office_document(self, filepath: str) -> Dict[str, Any]:
        """Extract text from Office documents"""
        try:
            # Try python-docx for Word files
            if filepath.endswith('.docx'):
                try:
                    from docx import Document
                    doc = Document(filepath)
                    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                    
                    return {
                        "method": "docx_python",
                        "extracted_text": text.strip(),
                        "confidence": 0.9,
                        "paragraphs": len(doc.paragraphs)
                    }
                except ImportError:
                    return {"method": "docx_no_lib", "error": "python-docx not available"}
            
            # Try openpyxl for Excel files
            elif filepath.endswith('.xlsx'):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(filepath)
                    text = ""
                    
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows():
                            for cell in row:
                                if cell.value:
                                    text += str(cell.value) + " "
                        text += "\n"
                    
                    return {
                        "method": "xlsx_openpyxl",
                        "extracted_text": text.strip(),
                        "confidence": 0.9,
                        "sheets": len(wb.sheetnames)
                    }
                except ImportError:
                    return {"method": "xlsx_no_lib", "error": "openpyxl not available"}
            
            # Try python-pptx for PowerPoint files
            elif filepath.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    prs = Presentation(filepath)
                    text = ""
                    
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + " "
                        text += "\n"
                    
                    return {
                        "method": "pptx_python",
                        "extracted_text": text.strip(),
                        "confidence": 0.9,
                        "slides": len(prs.slides)
                    }
                except ImportError:
                    return {"method": "pptx_no_lib", "error": "python-pptx not available"}
            
            else:
                return {"method": "office_unsupported", "error": "Unsupported office format"}
                
        except Exception as e:
            return {"method": "office_error", "error": str(e)}
    
    def extract_from_video(self, filepath: str) -> Dict[str, Any]:
        """Attempt to extract text from video (frames)"""
        try:
            # Try to extract frames and run OCR on them
            try:
                import cv2
                cap = cv2.VideoCapture(filepath)
                frame_count = 0
                all_text = ""
                confidences = []
                
                while frame_count < 10:  # Process first 10 frames
                    ret, frame = cap.read()
                    if not ret:
                        break
                    
                    # Convert to PIL Image
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_image = Image.fromarray(frame_rgb)
                    
                    # OCR on frame
                    text = pytesseract.image_to_string(pil_image, lang='eng')
                    if text.strip():
                        all_text += text + "\n"
                        
                        # Get confidence
                        data = pytesseract.image_to_data(pil_image, lang='eng', output_type=pytesseract.Output.DICT)
                        frame_confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
                        if frame_confidences:
                            confidences.extend(frame_confidences)
                    
                    frame_count += 1
                
                cap.release()
                
                avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                
                return {
                    "method": "video_frames_ocr",
                    "extracted_text": all_text.strip(),
                    "confidence": avg_confidence / 100,
                    "frames_processed": frame_count
                }
                
            except ImportError:
                return {"method": "video_no_opencv", "error": "OpenCV not available for video processing"}
                
        except Exception as e:
            return {"method": "video_error", "error": str(e)}
    
    def verify_ocr_accuracy(self, file_info: Dict[str, Any], ocr_result: Dict[str, Any]) -> Dict[str, Any]:
        """Verify OCR accuracy against expected content"""
        
        verification = {
            "file_info": file_info,
            "ocr_result": ocr_result,
            "expected_content": file_info.get("expected_content", []),
            "accuracy_score": 0.0,
            "content_matches": [],
            "missing_content": [],
            "extra_content": [],
            "is_accurate": False,
            "assessment": "unknown"
        }
        
        try:
            extracted_text = ocr_result.get("extracted_text", "").lower()
            expected_content = file_info.get("expected_content", [])
            
            # Check for expected content matches
            matches = []
            missing = []
            
            for expected in expected_content:
                if expected.lower() in extracted_text:
                    matches.append(expected)
                else:
                    missing.append(expected)
            
            verification["content_matches"] = matches
            verification["missing_content"] = missing
            
            # Calculate accuracy score
            if expected_content:
                verification["accuracy_score"] = len(matches) / len(expected_content)
            else:
                verification["accuracy_score"] = 0.0
            
            # Determine if OCR is accurate
            if verification["accuracy_score"] >= 0.7:
                verification["is_accurate"] = True
                verification["assessment"] = "ACCURATE"
            elif verification["accuracy_score"] >= 0.3:
                verification["is_accurate"] = False
                verification["assessment"] = "PARTIAL"
            else:
                verification["is_accurate"] = False
                verification["assessment"] = "INACCURATE"
            
            # Check for gibberish or nonsense
            if len(extracted_text.strip()) == 0:
                verification["assessment"] = "NO_TEXT_EXTRACTED"
            elif len(extracted_text.split()) < 3 and len(expected_content) > 0:
                verification["assessment"] = "INSUFFICIENT_TEXT"
            
        except Exception as e:
            verification["assessment"] = f"VERIFICATION_ERROR: {e}"
        
        return verification
    
    async def run_comprehensive_test(self):
        """Run comprehensive OCR testing"""
        
        # Download test content
        downloaded_files = await self.download_test_content()
        
        if not downloaded_files:
            print("❌ No files downloaded for testing")
            return
        
        print(f"\n🧪 RUNNING COMPREHENSIVE OCR TESTING")
        print("=" * 60)
        
        test_results = []
        
        for file_info in downloaded_files:
            print(f"\n🔍 Testing: {file_info['description']}")
            print(f"   File: {file_info['filename']}")
            print(f"   Category: {file_info['category']}")
            print(f"   Expected: {', '.join(file_info['expected_content'])}")
            
            # Extract text using OCR
            ocr_result = self.extract_text_from_file(file_info["filepath"], file_info["category"])
            
            # Verify accuracy
            verification = self.verify_ocr_accuracy(file_info, ocr_result)
            
            test_results.append(verification)
            
            # Display results
            if ocr_result.get("error"):
                print(f"   ❌ OCR Error: {ocr_result['error']}")
            else:
                method = ocr_result.get("method", "unknown")
                confidence = ocr_result.get("confidence", 0)
                text_length = len(ocr_result.get("extracted_text", ""))
                
                print(f"   ✅ OCR Method: {method}")
                print(f"   📊 Confidence: {confidence:.1%}")
                print(f"   📝 Text Length: {text_length} characters")
                
                # Show verification results
                assessment = verification["assessment"]
                accuracy = verification["accuracy_score"]
                
                print(f"   🎯 Assessment: {assessment}")
                print(f"   📈 Accuracy: {accuracy:.1%}")
                
                if verification["content_matches"]:
                    print(f"   ✅ Found: {', '.join(verification['content_matches'])}")
                
                if verification["missing_content"]:
                    print(f"   ❌ Missing: {', '.join(verification['missing_content'])}")
                
                # Show sample of extracted text
                extracted_sample = ocr_result.get("extracted_text", "")[:200]
                if extracted_sample:
                    print(f"   📄 Sample: \"{extracted_sample}...\"")
        
        # Generate comprehensive report
        self.generate_test_report(test_results)
        
        return test_results
    
    def generate_test_report(self, test_results: List[Dict[str, Any]]):
        """Generate comprehensive test report"""
        
        print(f"\n📊 COMPREHENSIVE OCR TEST REPORT")
        print("=" * 60)
        
        # Summary statistics
        total_tests = len(test_results)
        accurate_tests = sum(1 for r in test_results if r["is_accurate"])
        partial_tests = sum(1 for r in test_results if r["assessment"] == "PARTIAL")
        failed_tests = total_tests - accurate_tests - partial_tests
        
        accuracy_rate = accurate_tests / total_tests * 100 if total_tests > 0 else 0
        
        print(f"📈 SUMMARY STATISTICS")
        print(f"   Total Tests: {total_tests}")
        print(f"   Accurate: {accurate_tests} ({accuracy_rate:.1f}%)")
        print(f"   Partial: {partial_tests}")
        print(f"   Failed: {failed_tests}")
        
        # Category breakdown
        print(f"\n📂 CATEGORY PERFORMANCE")
        categories = {}
        for result in test_results:
            category = result["file_info"]["category"]
            if category not in categories:
                categories[category] = {"total": 0, "accurate": 0}
            categories[category]["total"] += 1
            if result["is_accurate"]:
                categories[category]["accurate"] += 1
        
        for category, stats in categories.items():
            rate = stats["accurate"] / stats["total"] * 100 if stats["total"] > 0 else 0
            print(f"   {category}: {stats['accurate']}/{stats['total']} ({rate:.1f}%)")
        
        # Method performance
        print(f"\n🔧 OCR METHOD PERFORMANCE")
        methods = {}
        for result in test_results:
            method = result["ocr_result"].get("method", "unknown")
            if method not in methods:
                methods[method] = {"total": 0, "accurate": 0}
            methods[method]["total"] += 1
            if result["is_accurate"]:
                methods[method]["accurate"] += 1
        
        for method, stats in methods.items():
            rate = stats["accurate"] / stats["total"] * 100 if stats["total"] > 0 else 0
            print(f"   {method}: {stats['accurate']}/{stats['total']} ({rate:.1f}%)")
        
        # Save detailed results
        report_data = {
            "timestamp": datetime.now().isoformat(),
            "summary": {
                "total_tests": total_tests,
                "accurate_tests": accurate_tests,
                "partial_tests": partial_tests,
                "failed_tests": failed_tests,
                "accuracy_rate": accuracy_rate
            },
            "category_performance": categories,
            "method_performance": methods,
            "detailed_results": test_results
        }
        
        with open('comprehensive_ocr_test_report.json', 'w') as f:
            json.dump(report_data, f, indent=2, default=str)
        
        print(f"\n📄 Detailed report saved to: comprehensive_ocr_test_report.json")
        
        # Overall assessment
        if accuracy_rate >= 80:
            print(f"\n🎉 OCR SYSTEM STATUS: EXCELLENT ✅")
        elif accuracy_rate >= 60:
            print(f"\n✅ OCR SYSTEM STATUS: GOOD ✅")
        elif accuracy_rate >= 40:
            print(f"\n⚠️  OCR SYSTEM STATUS: MODERATE ⚠️")
        else:
            print(f"\n❌ OCR SYSTEM STATUS: NEEDS IMPROVEMENT ❌")

async def main():
    """Main execution"""
    print("🧪 COMPREHENSIVE OCR TESTING SYSTEM")
    print("=" * 50)
    print("Testing OCR accuracy on diverse content types...")
    
    tester = ComprehensiveOCRTester()
    results = await tester.run_comprehensive_test()
    
    print(f"\n🏁 TESTING COMPLETE")
    print(f"Results: {len(results)} files tested")

if __name__ == "__main__":
    asyncio.run(main())
