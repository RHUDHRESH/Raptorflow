import io
from typing import Dict, List

from pptx import Presentation

from core.parsers.types import ParsedAsset


def _slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    return ""


def parse_pptx(content: bytes) -> ParsedAsset:
    presentation = Presentation(io.BytesIO(content))
    slide_entries: List[Dict[str, Any]] = []
    text_parts: List[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                slide_texts.append(text)
        slide_text = "\n".join(slide_texts).strip()
        if slide_text:
            text_parts.append(f"[Slide {index}]\n{slide_text}")
        title = _slide_title(slide)
        slide_entries.append(
            {
                "page_number": index,
                "title": title,
                "char_count": len(slide_text),
            }
        )

    combined_text = "\n\n".join(text_parts).strip()
    asset_metadata = {
        "page_count": len(slide_entries),
        "pages": slide_entries,
    }
    return ParsedAsset(text=combined_text, title="", metadata=asset_metadata)
